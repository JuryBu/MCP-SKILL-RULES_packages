import json
import os
import sys
import traceback
from typing import Any, Dict, List, Optional

from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx import Presentation


Rect = Dict[str, float]
SMALL_FONT_THRESHOLD_PT = 7.0
TITLE_TOP_VARIANCE_THRESHOLD_EMU = 180000.0
SIZE_VARIANCE_THRESHOLD_RATIO = 0.12
GAP_VARIANCE_THRESHOLD_RATIO = 0.35


def _rect_from_shape(shape: Any) -> Rect:
    return {
        "x0": float(shape.left),
        "y0": float(shape.top),
        "x1": float(shape.left + shape.width),
        "y1": float(shape.top + shape.height),
    }


def _rect_area(rect: Rect) -> float:
    return max(0.0, rect["x1"] - rect["x0"]) * max(0.0, rect["y1"] - rect["y0"])


def _overlap_area(a: Rect, b: Rect) -> float:
    width = min(a["x1"], b["x1"]) - max(a["x0"], b["x0"])
    height = min(a["y1"], b["y1"]) - max(a["y0"], b["y0"])
    return max(0.0, width) * max(0.0, height)


def _overlap_percent(a: Rect, b: Rect) -> float:
    smaller = min(_rect_area(a), _rect_area(b))
    if smaller <= 0:
        return 0.0
    return (_overlap_area(a, b) / smaller) * 100


def _shape_text(shape: Any) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    return (getattr(shape, "text", "") or "").strip()


def _shape_kind(shape: Any, text: str) -> str:
    if text:
        return "text"
    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
        return "image"
    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
        return "group"
    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.TABLE:
        return "table"
    return "shape"


def _font_size(shape: Any) -> Optional[float]:
    if not getattr(shape, "has_text_frame", False):
        return None
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.size:
                return float(run.font.size.pt)
    return None


def _text_clipping_sides(element: Dict[str, Any]) -> List[str]:
    text = str(element.get("text") or "")
    font_size = element.get("fontSize")
    if not text or not isinstance(font_size, (int, float)) or float(font_size) <= 0:
        return []

    bounds = element["bounds"]
    width = max(0.0, float(bounds["x1"] - bounds["x0"]))
    height = max(0.0, float(bounds["y1"] - bounds["y0"]))
    font_emu = float(font_size) * 12700.0
    lines = [line for line in text.splitlines() if line.strip()] or [text]
    longest_line = max(lines, key=len)
    estimated_line_width = len(longest_line) * font_emu * 0.55
    estimated_text_height = len(lines) * font_emu * 1.25
    sides: List[str] = []
    if estimated_line_width > width * 1.08:
        sides.append("right")
    if estimated_text_height > height * 1.08:
        sides.append("bottom")
    return sides


def _element_from_shape(shape: Any, z_order: int, page: int) -> Dict[str, Any]:
    text = _shape_text(shape)
    kind = _shape_kind(shape, text)
    element: Dict[str, Any] = {
        "type": kind,
        "name": getattr(shape, "name", f"shape-{z_order + 1}"),
        "text": text,
        "bounds": _rect_from_shape(shape),
        "zOrder": z_order,
        "source": "pptx",
        "page": page,
        "metadata": {
            "shapeType": str(getattr(shape, "shape_type", "unknown")),
        },
    }
    size = _font_size(shape)
    if size is not None:
        element["fontSize"] = size
    return element


def _union_rect(rects: List[Rect]) -> Rect:
    return {
        "x0": min(rect["x0"] for rect in rects),
        "y0": min(rect["y0"] for rect in rects),
        "x1": max(rect["x1"] for rect in rects),
        "y1": max(rect["y1"] for rect in rects),
    }


def _as_page_number(value: Any, default: Optional[int] = None) -> Optional[int]:
    if value is None:
        return default
    if isinstance(value, str) and value.lower() == "all":
        return None
    try:
        page = int(value)
    except (TypeError, ValueError) as exc:
        raise ValueError(f"page must be an integer, got {value!r}") from exc
    if page < 1:
        raise ValueError(f"page must be >= 1, got {page}")
    return page


def _open_pptx(pptx_path: str) -> Presentation:
    if not pptx_path:
        raise ValueError("pptxPath is required")
    if not os.path.exists(pptx_path):
        raise FileNotFoundError(f"PPTX file does not exist: {pptx_path}")
    return Presentation(pptx_path)


def _slide_indices(prs: Presentation, slide_num: Optional[int]) -> List[int]:
    if slide_num is None:
        return list(range(len(prs.slides)))
    slide_index = slide_num - 1
    if slide_index < 0 or slide_index >= len(prs.slides):
        raise IndexError(f"slide {slide_num} is out of range; PPTX has {len(prs.slides)} slide(s)")
    return [slide_index]


def extract_structure(pptx_path: str, slide_num: Optional[int] = None) -> List[Dict[str, Any]]:
    prs = _open_pptx(pptx_path)
    results: List[Dict[str, Any]] = []

    for slide_index in _slide_indices(prs, slide_num):
        slide = prs.slides[slide_index]
        page = slide_index + 1
        elements = [
            _element_from_shape(shape, z_order, page)
            for z_order, shape in enumerate(slide.shapes)
            if float(getattr(shape, "width", 0) or 0) > 0 and float(getattr(shape, "height", 0) or 0) > 0
        ]
        results.append({
            "page": page,
            "dimensions": {
                "width": float(prs.slide_width),
                "height": float(prs.slide_height),
                "unit": "EMU",
            },
            "elements": elements,
            "source": "pptx",
        })

    return results


def _overflow_sides(rect: Rect, width: float, height: float) -> List[str]:
    sides: List[str] = []
    if rect["x0"] < 0:
        sides.append("left")
    if rect["y0"] < 0:
        sides.append("top")
    if rect["x1"] > width:
        sides.append("right")
    if rect["y1"] > height:
        sides.append("bottom")
    return sides


def _severity_for_overlap(a: Dict[str, Any], b: Dict[str, Any]) -> str:
    types = {a["type"], b["type"]}
    if types == {"text"}:
        return "error"
    if "text" in types and "image" in types:
        text_elem = a if a["type"] == "text" else b
        image_elem = a if a["type"] == "image" else b
        if image_elem["zOrder"] > text_elem["zOrder"]:
            return "error"
        return "warning"
    if "text" in types:
        return "warning"
    return "info"


def detect_issues(
    pptx_path: str,
    slide_num: Optional[int] = None,
    checks: Optional[List[str]] = None,
    thresholds: Optional[Dict[str, Any]] = None,
) -> Dict[str, Any]:
    active_checks = set(checks or ["overlap", "overflow"])
    thresholds = thresholds or {}
    structures = extract_structure(pptx_path, slide_num)
    issues: List[Dict[str, Any]] = []

    for structure in structures:
        page = structure["page"]
        width = float(structure["dimensions"]["width"])
        height = float(structure["dimensions"]["height"])
        elements = structure["elements"]

        if "overflow" in active_checks:
            for element in elements:
                sides = _overflow_sides(element["bounds"], width, height)
                if not sides:
                    continue
                issues.append({
                    "type": "overflow",
                    "severity": "error",
                    "page": page,
                    "description": f"{element['name']} 超出幻灯片边界: {', '.join(sides)}",
                    "elements": [element],
                    "bounds": element["bounds"],
                    "metadata": {"sides": sides, "source": "pptx-native"},
                })

        if "overlap" in active_checks:
            for i, first in enumerate(elements):
                for second in elements[i + 1:]:
                    percent = _overlap_percent(first["bounds"], second["bounds"])
                    if percent < 2:
                        continue
                    issues.append({
                        "type": "overlap",
                        "severity": _severity_for_overlap(first, second),
                        "page": page,
                        "description": f"{first['name']} 与 {second['name']} 重叠 {percent:.1f}%",
                        "elements": [first, second],
                        "bounds": _union_rect([first["bounds"], second["bounds"]]),
                        "metadata": {
                            "overlapPercent": round(percent, 2),
                            "source": "pptx-native",
                        },
                    })

        if "readability" in active_checks:
            small_font_threshold = float(thresholds.get("smallFontPt") or SMALL_FONT_THRESHOLD_PT)
            for element in elements:
                if element.get("type") != "text":
                    continue
                font_size = element.get("fontSize")
                if isinstance(font_size, (int, float)) and float(font_size) < small_font_threshold:
                    issues.append({
                        "type": "small-font",
                        "severity": "warning",
                        "page": page,
                        "description": f"{element['name']} 字号 {float(font_size):.1f}pt 低于 {small_font_threshold:.1f}pt",
                        "elements": [element],
                        "bounds": element["bounds"],
                        "metadata": {
                            "check": "small-font",
                            "fontSize": float(font_size),
                            "threshold": small_font_threshold,
                            "source": "pptx-native",
                        },
                    })
                clipping_sides = _text_clipping_sides(element)
                if clipping_sides:
                    issues.append({
                        "type": "clipped",
                        "severity": "warning",
                        "page": page,
                        "description": f"{element['name']} 文本可能超出文本框: {', '.join(clipping_sides)}",
                        "elements": [element],
                        "bounds": element["bounds"],
                        "metadata": {
                            "check": "text-clipping",
                            "sides": clipping_sides,
                            "source": "pptx-native",
                        },
                    })

    if "alignment" in active_checks:
        issues.extend(_detect_title_alignment(
            structures,
            float(thresholds.get("titleTopVarianceEmu") or TITLE_TOP_VARIANCE_THRESHOLD_EMU),
        ))
        issues.extend(_detect_size_consistency(
            structures,
            float(thresholds.get("sizeVarianceRatio") or SIZE_VARIANCE_THRESHOLD_RATIO),
        ))
        issues.extend(_detect_uneven_spacing(
            structures,
            float(thresholds.get("gapVarianceRatio") or GAP_VARIANCE_THRESHOLD_RATIO),
        ))

    errors = sum(1 for issue in issues if issue["severity"] == "error")
    warnings = sum(1 for issue in issues if issue["severity"] == "warning")
    return {
        "summary": {
            "pages": len(structures),
            "elements": sum(len(structure["elements"]) for structure in structures),
            "issues": len(issues),
            "warnings": warnings,
            "errors": errors,
        },
        "issues": issues,
        "structure": structures,
    }


def _candidate_title(structure: Dict[str, Any]) -> Optional[Dict[str, Any]]:
    text_elements = [element for element in structure["elements"] if element.get("type") == "text" and element.get("text")]
    if not text_elements:
        return None
    named = [element for element in text_elements if "title" in str(element.get("name", "")).lower()]
    candidates = named or text_elements
    return sorted(
        candidates,
        key=lambda element: (
            float(element["bounds"]["y0"]),
            -float(element.get("fontSize") or 0),
            float(element["bounds"]["x0"]),
        ),
    )[0]


def _detect_title_alignment(structures: List[Dict[str, Any]], threshold: float) -> List[Dict[str, Any]]:
    titles = [title for title in (_candidate_title(structure) for structure in structures) if title]
    if len(titles) < 2:
        return []
    top_values = [float(title["bounds"]["y0"]) for title in titles]
    min_top = min(top_values)
    max_top = max(top_values)
    if max_top - min_top <= threshold:
        return []
    reference = sorted(top_values)[len(top_values) // 2]
    off_titles = [
        title for title in titles
        if abs(float(title["bounds"]["y0"]) - reference) > threshold
    ] or titles
    return [{
        "type": "misalignment",
        "severity": "warning",
        "page": int(off_titles[0].get("page") or 1),
        "description": f"跨页标题 top 位置差异 {max_top - min_top:.0f} EMU，超过阈值 {threshold:.0f} EMU",
        "elements": off_titles,
        "bounds": _union_rect([title["bounds"] for title in off_titles]),
        "metadata": {
            "check": "title-top-alignment",
            "minTop": min_top,
            "maxTop": max_top,
            "referenceTop": reference,
            "threshold": threshold,
            "source": "pptx-native",
        },
    }]


def _normalized_size_group(element: Dict[str, Any]) -> str:
    name = str(element.get("name") or "").lower()
    normalized = "".join(ch if ch.isalpha() else " " for ch in name).strip()
    first_word = normalized.split()[0] if normalized.split() else str(element.get("type") or "shape")
    return f"{element.get('type')}:{first_word}"


def _detect_size_consistency(structures: List[Dict[str, Any]], threshold: float) -> List[Dict[str, Any]]:
    groups: Dict[str, List[Dict[str, Any]]] = {}
    for structure in structures:
        for element in structure["elements"]:
            if element.get("type") not in {"text", "image", "shape"}:
                continue
            groups.setdefault(_normalized_size_group(element), []).append(element)

    issues: List[Dict[str, Any]] = []
    for group, elements in groups.items():
        if len(elements) < 2:
            continue
        widths = [float(element["bounds"]["x1"] - element["bounds"]["x0"]) for element in elements]
        heights = [float(element["bounds"]["y1"] - element["bounds"]["y0"]) for element in elements]
        avg_w = sum(widths) / len(widths)
        avg_h = sum(heights) / len(heights)
        width_ratio = ((max(widths) - min(widths)) / avg_w) if avg_w else 0
        height_ratio = ((max(heights) - min(heights)) / avg_h) if avg_h else 0
        if max(width_ratio, height_ratio) <= threshold:
            continue
        issues.append({
            "type": "inconsistent-size",
            "severity": "warning",
            "page": int(elements[0].get("page") or 1),
            "description": f"{group} 同类元素尺寸差异超过 {threshold:.0%}",
            "elements": elements,
            "bounds": _union_rect([element["bounds"] for element in elements]),
            "metadata": {
                "check": "size-consistency",
                "group": group,
                "widthVarianceRatio": round(width_ratio, 3),
                "heightVarianceRatio": round(height_ratio, 3),
                "threshold": threshold,
                "source": "pptx-native",
            },
        })
    return issues


def _detect_uneven_spacing(structures: List[Dict[str, Any]], threshold: float) -> List[Dict[str, Any]]:
    issues: List[Dict[str, Any]] = []
    for structure in structures:
        text_elements = [
            element for element in structure["elements"]
            if element.get("type") == "text" and float(element["bounds"]["x1"] - element["bounds"]["x0"]) > 0
        ]
        rows: Dict[int, List[Dict[str, Any]]] = {}
        row_bucket = 120000
        for element in text_elements:
            bucket = round(float(element["bounds"]["y0"]) / row_bucket)
            rows.setdefault(bucket, []).append(element)
        for bucket, row in rows.items():
            if len(row) < 4:
                continue
            ordered = sorted(row, key=lambda element: float(element["bounds"]["x0"]))
            gaps = [
                float(right["bounds"]["x0"] - left["bounds"]["x1"])
                for left, right in zip(ordered, ordered[1:])
                if right["bounds"]["x0"] >= left["bounds"]["x1"]
            ]
            if len(gaps) < 3:
                continue
            avg_gap = sum(gaps) / len(gaps)
            if avg_gap <= 0:
                continue
            variance = (max(gaps) - min(gaps)) / avg_gap
            if variance <= threshold:
                continue
            issues.append({
                "type": "uneven-spacing",
                "severity": "info",
                "page": int(structure["page"]),
                "description": f"第 {structure['page']} 页同一行元素间距差异 {variance:.2f}，超过阈值 {threshold:.2f}",
                "elements": ordered,
                "bounds": _union_rect([element["bounds"] for element in ordered]),
                "metadata": {
                    "check": "row-gap-consistency",
                    "rowBucket": bucket,
                    "gaps": [round(gap, 2) for gap in gaps],
                    "varianceRatio": round(variance, 3),
                    "threshold": threshold,
                    "source": "pptx-native",
                },
            })
    return issues


def search_text_region(pptx_path: str, target: str, slide_num: Optional[int] = None) -> Dict[str, Any]:
    if not target or not str(target).strip():
        raise ValueError("target is required")

    prs = _open_pptx(pptx_path)
    target_text = str(target).strip().lower()
    slide_indices = _slide_indices(prs, slide_num)

    for slide_index in slide_indices:
        if slide_index < 0 or slide_index >= len(prs.slides):
            raise IndexError(f"slide {slide_index + 1} is out of range; PPTX has {len(prs.slides)} slide(s)")
        slide = prs.slides[slide_index]
        matches: List[Rect] = []
        names: List[str] = []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text = (getattr(shape, "text", "") or "").strip()
            if target_text not in text.lower():
                continue
            matches.append(_rect_from_shape(shape))
            names.append(getattr(shape, "name", f"shape-{len(names) + 1}"))

        if matches:
            return {
                "found": True,
                "target": target,
                "page": slide_index + 1,
                "rect": _union_rect(matches),
                "matches": len(matches),
                "elements": names,
                "dimensions": {
                    "width": float(prs.slide_width),
                    "height": float(prs.slide_height),
                    "unit": "EMU",
                },
            }

    return {
        "found": False,
        "target": target,
        "page": slide_num,
        "rect": None,
        "matches": 0,
        "elements": [],
        "dimensions": {
            "width": float(prs.slide_width),
            "height": float(prs.slide_height),
            "unit": "EMU",
        },
    }


def _params(payload: Dict[str, Any]) -> Dict[str, Any]:
    params = dict(payload.get("params") or {})
    for key, value in payload.items():
        if key not in {"action", "params"} and key not in params:
            params[key] = value
    return params


def _dispatch(payload: Dict[str, Any]) -> Any:
    action = payload.get("action")
    params = _params(payload)
    pptx_path = params.get("pptxPath") or params.get("pptx_path") or params.get("path")
    page_num = _as_page_number(params.get("page") or params.get("pageNum") or params.get("page_num"))

    if action == "search_text_region":
        return search_text_region(pptx_path, str(params.get("target") or ""), page_num)
    if action == "extract_structure":
        return extract_structure(pptx_path, page_num)
    if action == "detect_issues":
        return detect_issues(
            pptx_path,
            page_num,
            list(params.get("checks") or ["overlap", "overflow"]),
            params.get("thresholds"),
        )

    raise ValueError(f"unknown action: {action!r}")


def main() -> int:
    try:
        raw = sys.stdin.read()
        if not raw.strip():
            raise ValueError("stdin JSON payload is required")
        payload = json.loads(raw)
        if not isinstance(payload, dict):
            raise ValueError("stdin JSON payload must be an object")
        result = _dispatch(payload)
        print(json.dumps({"ok": True, "action": payload.get("action"), "result": result}, ensure_ascii=False))
        return 0
    except Exception as exc:
        print(json.dumps({
            "ok": False,
            "error": {
                "type": exc.__class__.__name__,
                "message": str(exc),
                "traceback": traceback.format_exc(),
            },
        }, ensure_ascii=False))
        return 1


if __name__ == "__main__":
    raise SystemExit(main())
